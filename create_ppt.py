"""
ErgoGuard Chair - PowerPoint Generator
Creates a 10-slide presentation PPT file
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Define colors
    primary_color = RGBColor(102, 126, 234)  # Purple/Blue
    accent_green = RGBColor(76, 175, 80)
    accent_amber = RGBColor(255, 152, 0)
    accent_red = RGBColor(244, 67, 54)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "ErgoGuard Chair"
    subtitle1.text = "Context-Aware Posture + Stress + Fatigue Prevention System\n\n[Your Name]\n[College Name]\nComputer Science Engineering\nProject Proposal\n[Date]"
    
    # Slide 2: Problem Statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    title2.text = "The Problem We're Solving"
    content2.text = (
        "• 70% of IT workers suffer from chronic back pain\n"
        "• Average sitting time: 9.3 hours/day\n"
        "• Existing solutions fail:\n"
        "  - Apple Watch: Can't measure posture\n"
        "  - Camera-based: Privacy concerns\n"
        "  - Simple reminders: Not effective\n"
        "• Healthcare costs: ₹50,000+ per employee/year\n"
        "• Productivity loss: 15-20%"
    )
    
    # Slide 3: Proposed Solution
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    title3.text = "ErgoGuard Chair - Our Solution"
    content3.text = (
        "Hardware + Software Product\n\n"
        "• Real-time posture detection using pressure sensors\n"
        "• Sedentary behavior monitoring via micro-movements\n"
        "• Adaptive interventions (vibration, LED feedback)\n"
        "• Web dashboard for analytics\n\n"
        "Key Innovation:\n"
        "• Measures what wearables can't (posture, weight distribution)\n"
        "• Provides physical interventions\n"
        "• Privacy-preserving design"
    )
    
    # Slide 4: System Architecture
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    title4.text = "Proposed System Architecture"
    content4.text = (
        "Sensors → ESP32 → Backend → Frontend\n\n"
        "Hardware Layer:\n"
        "• 8x FSR Pressure Sensors (seat + backrest)\n"
        "• MPU6050 IMU (motion detection)\n"
        "• MQ135 + DHT22 (environment)\n"
        "• Vibration Motors + RGB LED\n\n"
        "Software Layer:\n"
        "• ESP32 Firmware (Arduino)\n"
        "• Node.js Backend + MongoDB\n"
        "• React Dashboard"
    )
    
    # Slide 5: Core Algorithm
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    title5.text = "Ergo-Stress Score Algorithm"
    content5.text = (
        "Base Score = 100\n\n"
        "Posture Penalty = f(deviation)\n"
        "  - Forward lean: -20\n"
        "  - Side lean: -15\n"
        "  - Slouching: -25\n\n"
        "Asymmetry Penalty = f(imbalance)\n"
        "  - Left-right: -10\n"
        "  - Front-back: -5\n\n"
        "Sedentary Penalty = f(movement, duration)\n"
        "  - No movement 30min: -15\n"
        "  - No movement 60min: -30\n\n"
        "Final Score = Base - (Posture + Asymmetry + Sedentary)\n\n"
        "Novelty: Multi-modal sensor fusion"
    )
    
    # Slide 6: Key Features
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    title6.text = "Proposed Features"
    content6.text = (
        "1. Real-Time Monitoring:\n"
        "• Posture Score (0-100)\n"
        "• Asymmetry Score\n"
        "• Sedentary Score\n"
        "• Ergo-Stress Score\n\n"
        "2. Adaptive Interventions:\n"
        "• Visual: RGB LED (green/amber/red)\n"
        "• Haptic: Vibration motors\n"
        "• Smart alerts & recommendations\n\n"
        "3. Analytics Dashboard:\n"
        "• Real-time visualization\n"
        "• Historical trends\n"
        "• Posture breakdown"
    )
    
    # Slide 7: Implementation Plan
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    title7.text = "2-Month Development Timeline"
    content7.text = (
        "Week 1-2: Hardware Setup\n"
        "• Component procurement\n"
        "• Sensor assembly\n"
        "• ESP32 integration\n\n"
        "Week 3-4: Firmware Development\n"
        "• Sensor reading\n"
        "• Algorithm implementation\n"
        "• WiFi communication\n\n"
        "Week 5-6: Backend Development\n"
        "• Node.js server\n"
        "• MongoDB database\n"
        "• API endpoints\n\n"
        "Week 7: Frontend Development\n"
        "• React dashboard\n"
        "• Real-time visualization\n\n"
        "Week 8: Integration & Testing\n"
        "• End-to-end testing\n"
        "• Demo preparation"
    )
    
    # Slide 8: Innovation & Patent Potential
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    title8.text = "Patent-Worthy Innovations"
    content8.text = (
        "5 Key Novelty Points:\n\n"
        "1. Multi-Modal Sensor Fusion\n"
        "   - Pressure + IMU + Environment → Unified score\n\n"
        "2. Adaptive Intervention Policy\n"
        "   - Context-aware, not fixed intervals\n\n"
        "3. Privacy-Preserving Work Integration\n"
        "   - Activity patterns without content monitoring\n\n"
        "4. Auto-Calibration Method\n"
        "   - Adapts to individual body types\n\n"
        "5. Sedentary Risk Index\n"
        "   - Micro-movement analysis algorithm"
    )
    
    # Slide 9: Expected Impact & Market
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    title9.text = "Impact & Commercial Potential"
    content9.text = (
        "Health Benefits:\n"
        "• Reduced back/neck pain\n"
        "• Improved posture awareness\n"
        "• Better break habits\n\n"
        "Market Potential:\n"
        "• B2B: IT companies, call centers\n"
        "• B2C: Home office workers, students\n"
        "• Market Size: Millions of IT workers\n\n"
        "Cost Analysis:\n"
        "• Prototype: ₹3,200-4,000\n"
        "• Production: ₹2,000-2,500/unit\n"
        "• Clear ROI for companies"
    )
    
    # Slide 10: Conclusion
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    title10.text = "ErgoGuard Chair - Proposed Solution"
    content10.text = (
        "Summary:\n"
        "✅ Solves Real Problem\n"
        "✅ Patent-Worthy Innovation\n"
        "✅ Complete Solution Design\n"
        "✅ Feasible Implementation\n"
        "✅ Clear Market Value\n\n"
        "Next Steps:\n"
        "1. Component procurement\n"
        "2. Hardware assembly\n"
        "3. Development (2 months)\n"
        "4. Testing & demo\n\n"
        "Call to Action:\n"
        "\"Making workplace health monitoring accessible, intelligent, and effective.\""
    )
    
    # Apply formatting to all slides
    for slide in prs.slides:
        # Format titles
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    # Title formatting
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(44)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = primary_color
                else:
                    # Content formatting
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(20)
                        paragraph.space_after = Pt(12)
    
    # Save presentation
    output_path = r"d:\projects\ErgoGuard-Chair\ErgoGuard_Chair_Proposal_10_Slides.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint created successfully!")
    print(f"📁 Location: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    try:
        output_file = create_presentation()
        print("\n🎉 Presentation ready! Open it in PowerPoint.")
    except ImportError:
        print("❌ Error: python-pptx library not installed")
        print("📦 Installing python-pptx...")
        import subprocess
        import sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        print("✅ Installation complete! Run the script again.")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

